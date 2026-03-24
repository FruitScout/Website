"""
FarmAgent Agave Sales Deck — PowerPoint Generator
Generates FarmAgent-Agave.pptx (21 slides, 16:9 widescreen)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Constants ──────────────────────────────────────────────────────────────
W, H = Inches(13.333), Inches(7.5)   # 16:9 widescreen

# Brand colours
VOID       = RGBColor(0x08, 0x0E, 0x14)
SURFACE    = RGBColor(0x0D, 0x1A, 0x26)
SURFACE2   = RGBColor(0x13, 0x22, 0x33)
BORDER     = RGBColor(0x1E, 0x33, 0x47)
BLUE       = RGBColor(0x00, 0x6E, 0xA7)
BLUE_B     = RGBColor(0x0A, 0x8F, 0xD4)
GREEN      = RGBColor(0xA4, 0xCA, 0x67)
GOLD       = RGBColor(0xD2, 0xA8, 0x44)
TEXT       = RGBColor(0xEA, 0xF2, 0xF8)
TEXT2      = RGBColor(0xA8, 0xBD, 0xD0)
GREY_M     = RGBColor(0x4A, 0x62, 0x78)
GREY_L     = RGBColor(0xA8, 0xBD, 0xD0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

IMG_DIR = os.path.join(os.path.dirname(__file__), "..", "images")
def img(name): return os.path.join(IMG_DIR, name)


# ── Helpers ─────────────────────────────────────────────────────────────────

def set_bg(slide, color=None):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or VOID

def add_bg_image(slide, path, brightness=1.0):
    """Add a full-bleed background image, then dark overlay."""
    if not os.path.exists(path):
        set_bg(slide)
        return
    pic = slide.shapes.add_picture(path, 0, 0, W, H)
    pic.zorder = 0
    # Dark overlay rectangle
    alpha = int((1 - brightness) * 200)   # rough opacity
    overlay = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, H
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = VOID
    overlay.line.fill.background()
    # Make it semi-transparent via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    sp_tree = overlay._element
    solid_fill = sp_tree.find('.//' + qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is None:
            srgb = solid_fill.find(qn('a:sysClr'))
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(int(alpha * 1000)))
    overlay.zorder = 1

def txb(slide, left, top, width, height, text, size=Pt(14), bold=False,
        color=None, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    """Add a text box and return the text frame."""
    if color is None:
        color = TEXT
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tf

def add_rule(slide, left, top, width=Inches(0.6), color=None):
    """Thin horizontal rule."""
    line = slide.shapes.add_shape(1, left, top, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color or GOLD
    line.line.fill.background()

def add_tag(slide, left, top, text, addon=False):
    """Small capsule tag label."""
    color = GOLD if addon else GREEN
    tf = txb(slide, left, top, Inches(5), Inches(0.35),
             text.upper(), size=Pt(9), color=color, font="Calibri")
    return tf

def add_pill(slide, left, top, width, height, bg_color, text, text_color=WHITE):
    """Filled pill / callout box."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = bg_color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return tf

def footer_bar(slide):
    """Thin gradient-style footer line (solid blue approximation)."""
    bar = slide.shapes.add_shape(1, 0, H - Pt(3), W, Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

def bullet_list(slide, left, top, width, height, items, size=Pt(11.5),
                bullet_color=GREEN, text_color=TEXT2, gap=Inches(0.05)):
    """Add a bulleted list. items: list of (bold_prefix, rest_text) tuples or plain strings."""
    bx = slide.shapes.add_textbox(left, top, width, height)
    tf = bx.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        # Bullet character
        br = p.add_run()
        br.text = "■  "
        br.font.size = Pt(8)
        br.font.color.rgb = bullet_color
        br.font.name = "Calibri"
        if isinstance(item, tuple):
            bold_part, rest = item
            r2 = p.add_run()
            r2.text = bold_part
            r2.font.bold = True
            r2.font.size = size
            r2.font.color.rgb = TEXT
            r2.font.name = "Calibri"
            r3 = p.add_run()
            r3.text = rest
            r3.font.size = size
            r3.font.color.rgb = text_color
            r3.font.name = "Calibri"
        else:
            r2 = p.add_run()
            r2.text = str(item)
            r2.font.size = size
            r2.font.color.rgb = text_color
            r2.font.name = "Calibri"
    return tf

def add_roi_box(slide, left, top, width, height, text, border_color=None):
    """Left-border callout box."""
    bc = border_color or GREEN
    # Border bar
    bar = slide.shapes.add_shape(1, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = bc
    bar.line.fill.background()
    # Content box
    box = slide.shapes.add_shape(1, left + Pt(8), top, width - Pt(8), height)
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    return tf

def right_image(slide, fname, left, top, width, height):
    """Add image if file exists, else placeholder rect."""
    path = img(fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        r = slide.shapes.add_shape(1, left, top, width, height)
        r.fill.solid()
        r.fill.fore_color.rgb = SURFACE2
        r.line.color.rgb = BORDER

# ── Slide builders ───────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_image(sl, img("deck_aerial_dark.png"), brightness=0.38)
    set_bg(sl, VOID)
    # Logo
    txb(sl, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
        "FRUITSCOUT", size=Pt(12), bold=True, color=BLUE_B)
    # Tag
    txb(sl, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4),
        "THE INTELLIGENT FARM OS  ·  AGAVE EDITION",
        size=Pt(10), color=GREEN)
    # Headline
    box = sl.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(8), Inches(2.8))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [("Seven years.", TEXT), ("Managed blind.", TEXT), ("Not anymore.", GREEN)]
    first = True
    for text, col in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.size = Pt(52)
        run.font.name = "Calibri"
        run.font.bold = False
        run.font.color.rgb = col
    # Sub
    txb(sl, Inches(0.8), Inches(4.7), Inches(7.5), Inches(1.2),
        "FarmAgent gives agave producers the eyes, the intelligence, and the operating system "
        "to see exactly what is happening across their estate — and act on it before it costs them a harvest.",
        size=Pt(13), color=TEXT2)
    # Confidential
    txb(sl, Inches(0.8), Inches(6.9), Inches(6), Inches(0.4),
        "FruitScout · FarmAgent · Confidential", size=Pt(9), color=GREY_M)
    footer_bar(sl)


def slide_divider(prs, section_num, section_label, headline_lines, bg_img="deck_aerial_dark.png", brightness=0.30):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(sl, img(bg_img), brightness=brightness)
    # Ghost section number
    txb(sl, Inches(9.5), Inches(0.2), Inches(3.5), Inches(4),
        section_num, size=Pt(180), color=RGBColor(0x1E, 0x33, 0x47),
        font="Calibri")
    # Section label
    txb(sl, Inches(0.8), Inches(4.3), Inches(9), Inches(0.4),
        section_label.upper(), size=Pt(10), color=GREY_M)
    # Headline
    box = sl.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(9.5), Inches(2.3))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, col in headline_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.size = Pt(46)
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = col
    footer_bar(sl)


def slide_two_col(prs, tag, headline, left_content_fn, right_img_name, addon=False, rule_color=None):
    """Generic two-column content slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, VOID)
    add_tag(sl, Inches(0.7), Inches(0.35), tag, addon=addon)
    # Headline
    txb(sl, Inches(0.7), Inches(0.75), Inches(12.5), Inches(0.85),
        headline, size=Pt(22), bold=True, color=TEXT)
    add_rule(sl, Inches(0.7), Inches(1.65), color=rule_color or GOLD)
    # Left column content
    left_content_fn(sl)
    # Right image
    right_image(sl, right_img_name, Inches(7.0), Inches(1.0), Inches(5.9), Inches(5.8))
    footer_bar(sl)
    return sl


def slide_full_content(prs, tag, headline, content_fn, addon=False):
    """Full-width content slide (no right image)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, VOID)
    add_tag(sl, Inches(0.7), Inches(0.35), tag, addon=addon)
    txb(sl, Inches(0.7), Inches(0.75), Inches(12.5), Inches(0.85),
        headline, size=Pt(22), bold=True, color=TEXT)
    add_rule(sl, Inches(0.7), Inches(1.65))
    content_fn(sl)
    footer_bar(sl)
    return sl


# ── Individual slide content functions ───────────────────────────────────────

def content_problem(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.9),
        "A Blue Weber agave estate is a 7-year illiquid investment. Every undetected mortality event, "
        "premature harvest, or poorly executed Jima cut compounds in silence — often for years before the damage appears on the P&L.",
        size=Pt(11), color=TEXT2)
    bullet_list(sl, Inches(0.7), Inches(2.85), Inches(6.0), Inches(2.8),
        [("Mortality Bleed", " — Paper ledgers claim 40,000 plants. A drone survey reveals 36,200. Growing silently since Year 2."),
         ("No production visibility", " — Monthly biomass accumulation is invisible. Distillery contracts signed on gut-feel."),
         ("Jima quality variance", " — Cut quality determines ART sugar content. Most operations have no system to measure it."),
         ("Reactive management", " — Problems discovered on walkabouts. Intervention window closes before action is taken.")])
    # Stat callout
    cb = add_roi_box(sl, Inches(0.7), Inches(5.8), Inches(6.0), Inches(1.3), "", GOLD)
    p = cb.paragraphs[0]
    run = p.add_run()
    run.text = "$148,500 "
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = GOLD
    run.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = "— unrecovered revenue from 10% undetected mortality on 100 ha over 5 years"
    r2.font.size = Pt(11)
    r2.font.color.rgb = TEXT2
    r2.font.name = "Calibri"


def content_overview(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.8),
        "FarmAgent replaces operational blindness with complete visibility — from 120 metres in the air "
        "to the screen in a Jima worker's hand.",
        size=Pt(11.5), color=TEXT2)
    bullet_list(sl, Inches(0.7), Inches(2.7), Inches(6.0), Inches(2.5),
        [("Aerial Sight", " — Computer vision on drones counts every living plant, maps every gap."),
         ("Ground Sight", " — Smartphones measure production, size, and location with instrument-grade precision."),
         ("Operational Sight", " — The Central Brain sees goals, resources, processes, and plans simultaneously."),
         ("Longitudinal Sight", " — Farm Brain sees across seasons, surfacing patterns no single year reveals.")])
    roi = add_roi_box(sl, Inches(0.7), Inches(5.5), Inches(6.0), Inches(0.8), "", GREEN)
    p = roi.paragraphs[0]
    r = p.add_run()
    r.text = "FarmAgent turns a 7-year blind wait into a deterministic, auditable supply chain."
    r.font.size = Pt(11.5)
    r.font.bold = True
    r.font.color.rgb = GREEN
    r.font.name = "Calibri"


def content_aerial(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.4),
        "Advanced Inventory", size=Pt(12), bold=True, color=TEXT)
    bullet_list(sl, Inches(0.7), Inches(2.3), Inches(6.0), Inches(2.1),
        [("Object recognition", " identifies every living plant and cross-references against the prior year's spatial map."),
         ("Mortality Bleed Rate", " — the exact gap between what the ledger says exists and what is biologically alive."),
         ("Every missing plant", " GPS-tagged and queued as a replanting opportunity — no manual flagging."),
         ("Year-over-year canopy comparison", " flags declining vigor before losses are irreversible.")])
    txb(sl, Inches(0.7), Inches(4.5), Inches(6.0), Inches(0.4),
        "Targeted Scouting", size=Pt(12), bold=True, color=TEXT)
    bullet_list(sl, Inches(0.7), Inches(4.95), Inches(6.0), Inches(1.2),
        ["AI generates GPS routes sending scouts only to locations that need attention. Blanket scouting eliminated.",
         ("40–60% scouting labor reduction", " on large estates.")])


def content_field_agent(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.4),
        "Guided Work Tickets", size=Pt(12), bold=True, color=TEXT)
    bullet_list(sl, Inches(0.7), Inches(2.3), Inches(6.0), Inches(1.5),
        ["Every task dispatched with GPS block boundary, step sequence, inputs, timing, and quality criteria.",
         "Large instruction cards, photo prompts, confirmation checkpoints — no training required.",
         ("Offline-first", " — full task syncs before the shift. Works with zero signal.")])
    txb(sl, Inches(0.7), Inches(3.95), Inches(6.0), Inches(0.4),
        "Computer Vision Built Into the Workflow", size=Pt(12), bold=True, color=TEXT)
    bullet_list(sl, Inches(0.7), Inches(4.4), Inches(6.0), Inches(2.2),
        [("Jima crew:", " photograph the cut → computer vision grades A, B, or C before the worker moves on."),
         ("Field worker:", " photograph a rain-flooded road blocking Block 4 → FarmAgent reroutes the harvest convoy before a single truck leaves the yard."),
         ("Spray crew:", " photograph the team at task start → PPE class verified before the task unlocks.")])


def content_measurement(sl):
    items = [
        ("Visual Scale — Monthly Production\n",
         "AI converts smartphone imagery into volumetric piña estimates — no weighbridge. "
         "Updated monthly, block by block: the heartbeat of the business."),
        ("Visual Calipers — Precision Sizing\n",
         "Reference card establishes scale in every photo. AI returns mm-accurate readings in <10 sec. "
         "40–60 plants/hour vs. 8–12 with physical calipers."),
        ("Visual Localizer — GPS Anomaly Tagging\n",
         "Tag a diseased plant or mortality gap in Field Agent. "
         "Pin appears on every crew member's map within seconds.")
    ]
    top = Inches(1.85)
    for bold, rest in items:
        txb(sl, Inches(0.7), top, Inches(6.0), Inches(0.35),
            bold.strip(), size=Pt(12), bold=True, color=TEXT)
        top += Inches(0.35)
        txb(sl, Inches(0.7), top, Inches(6.0), Inches(0.7),
            rest, size=Pt(11), color=TEXT2)
        top += Inches(0.75)
        line = sl.shapes.add_shape(1, Inches(0.7), top, Inches(5.8), Pt(1))
        line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
        top += Inches(0.08)


def content_central_brain(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(12.5), Inches(0.65),
        "Every AI recommendation is bounded by four real-world pillars before it is actioned. "
        "The AI cannot prescribe resources that don't exist, violate approved processes, or act against your stated financial goals.",
        size=Pt(11), color=TEXT2)
    # Table header
    cols = ["Pillar", "Question it answers", "What it governs"]
    col_widths = [Inches(2.6), Inches(4.0), Inches(5.5)]
    col_starts = [Inches(0.7), Inches(3.3), Inches(7.3)]
    rows = [
        ("Goal Management", "What are we trying to achieve?", "Financial targets, tonnage KPIs, quality thresholds"),
        ("Resource Management", "What do we have?", "Labor, equipment, biological inventory"),
        ("Process Management", "How should it be done?", "SOPs, labor rates, chemical dosages, timings"),
        ("Plan Management", "When and in what sequence?", "The 7-year agave lifecycle, dynamically sequenced"),
    ]
    top = Inches(2.6)
    for c, ci, cw in zip(cols, col_starts, col_widths):
        txb(sl, ci, top, cw, Inches(0.35), c.upper(),
            size=Pt(9), color=GREY_M, bold=True)
    top += Inches(0.35)
    line = sl.shapes.add_shape(1, Inches(0.7), top, Inches(12.2), Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
    top += Inches(0.08)
    for row in rows:
        for val, ci, cw, bold in zip(row, col_starts, col_widths, [True, False, False]):
            txb(sl, ci, top, cw, Inches(0.5), val,
                size=Pt(11), color=TEXT if bold else TEXT2, bold=bold)
        top += Inches(0.5)
        line2 = sl.shapes.add_shape(1, Inches(0.7), top, Inches(12.2), Pt(1))
        line2.fill.solid(); line2.fill.fore_color.rgb = SURFACE2; line2.line.fill.background()
        top += Inches(0.06)
    roi = add_roi_box(sl, Inches(0.7), Inches(6.45), Inches(12.2), Inches(0.7), "", GREEN)
    p = roi.paragraphs[0]
    r = p.add_run()
    r.text = "The Guardrail Engine: this is what makes FarmAgent operationally trustworthy — not just impressive in a demo."
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"


def content_seven_year(sl):
    # Timeline bar
    bar_top = Inches(2.3)
    bar = sl.shapes.add_shape(1, Inches(0.7), bar_top + Inches(0.15), Inches(12.0), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    phases = [
        ("Yr 1–2", "Planting &\nEstablishment", False),
        ("Yr 2–3", "Monthly\nProduction Pulse", False),
        ("Yr 3–4", "Cohort\nSequencing", False),
        ("Yr 4–5", "Supply Contract\nForecasting", False),
        ("Yr 5–6", "Jima Readiness\nAssessment", False),
        ("Yr 6–7", "Terminal Jima &\nReplanting Queue", True),  # Gold = harvest
    ]
    for i, (yr, phase, is_harvest) in enumerate(phases):
        x = Inches(0.7) + Inches(2.0) * i
        dot = sl.shapes.add_shape(9, x + Inches(0.8), bar_top + Inches(0.02), Inches(0.26), Inches(0.26))
        dot.fill.solid(); dot.fill.fore_color.rgb = GOLD if is_harvest else GREEN
        dot.line.fill.background()
        txb(sl, x, bar_top + Inches(0.45), Inches(1.9), Inches(0.35),
            yr, size=Pt(10), bold=True, color=TEXT2, align=PP_ALIGN.CENTER)
        txb(sl, x, bar_top + Inches(0.82), Inches(1.9), Inches(0.7),
            phase, size=Pt(9), color=GREY_M, align=PP_ALIGN.CENTER)
    bullet_list(sl, Inches(0.7), Inches(3.9), Inches(12.2), Inches(1.8),
        [("Cohort sequencing", " — Slow-growth blocks harvested first; fast-compounding blocks keep accumulating."),
         ("Dynamic replanning", " — Mortality event or contract change triggers automatic plan revision in real time."),
         ("Replanting windows", " — Year 1–2 mortality: same cohort. Year 5: permanent loss. FarmAgent flags the difference.")])
    roi = add_roi_box(sl, Inches(0.7), Inches(5.8), Inches(12.2), Inches(0.9), "", BLUE_B)
    p = roi.paragraphs[0]
    r = p.add_run()
    r.text = '"Most agave operations make a new plan every year. FarmAgent runs one plan across 7 years — and updates it every time the biology changes."'
    r.font.size = Pt(11.5); r.font.italic = True; r.font.color.rgb = TEXT2; r.font.name = "Calibri"


def content_jima(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.65),
        "An A-Grade Jima cut delivers 15–22% higher ART sugar content than a C-Grade cut. "
        "FarmAgent grades every cut at full workforce scale — built into the Jima work ticket in the Field Agent app.",
        size=Pt(11.5), color=TEXT2)
    steps = [
        "Field Agent prompts the Jima worker to photograph the completed cut.",
        "Computer vision grades A, B, or C — cut angle, stub length, leaf trim, piña integrity — before the worker moves on.",
        "Sub-standard cuts flag instantly in the Manager Hub QA Inbox — sorted by risk.",
        "One tap → rework ticket dispatched to crew leader with photo attached. Crew stays until resolved.",
    ]
    top = Inches(2.65)
    for i, step in enumerate(steps, 1):
        # Step number
        dot = sl.shapes.add_shape(9, Inches(0.7), top, Inches(0.28), Inches(0.28))
        dot.fill.solid(); dot.fill.fore_color.rgb = BLUE; dot.line.fill.background()
        txb(sl, Inches(1.1), top, Inches(5.5), Inches(0.45),
            step, size=Pt(11), color=TEXT2)
        top += Inches(0.52)
    roi = add_roi_box(sl, Inches(0.7), Inches(5.55), Inches(6.0), Inches(1.1), "", GREEN)
    p = roi.paragraphs[0]
    r = p.add_run(); r.text = "Moving 100 ha from 55% → 82% A-Grade:\n"
    r.font.size = Pt(11.5); r.font.color.rgb = TEXT2; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = "$68,000–$113,000 USD"
    r2.font.size = Pt(20); r2.font.bold = True; r2.font.color.rgb = GOLD; r2.font.name = "Calibri"
    r3 = p.add_run(); r3.text = " recovered per harvest season"
    r3.font.size = Pt(11.5); r3.font.color.rgb = TEXT2; r3.font.name = "Calibri"


def content_farm_brain(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.7),
        "Where the Central Brain gives sight of today, Farm Brain gives sight across time — absorbing agronomic expertise, "
        "historical telemetry, and commercial intelligence into a continuously improving reasoning engine.",
        size=Pt(11.5), color=TEXT2)
    bullet_list(sl, Inches(0.7), Inches(2.7), Inches(6.0), Inches(2.5),
        [("Agave Intelligence Corpus", " — Blue Weber growth curves, ART accumulation models, weevil pressure cycles, mortality benchmarks."),
         ("Continual Improvement", " — Lean/Kaizen applied to operations. Surfaces OPEX reductions from your own estate's data."),
         ("Business Intelligence", " — COGS modeling, distillery contract structuring, multi-cohort capital planning.")])
    txb(sl, Inches(0.7), Inches(5.35), Inches(6.0), Inches(0.9),
        "By Year 3, the system is reasoning with 3 years of your estate's own telemetry. "
        "This is what makes it progressively harder to displace.",
        size=Pt(10), color=GREY_M, italic=True)


def content_personal_agent(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.65),
        "Every drone count, tonnage forecast, Jima quality score, and plan milestone — accessible through a single conversation. "
        "In Spanish or English, with your actual field numbers behind every answer.",
        size=Pt(11.5), color=TEXT2)
    chats = [
        ('"How is Block 7 tracking against the 7-year plan?"',
         "Block 7 is in Year 4, tracking 11% behind the biomass curve. Candidate for early harvest sequencing. Want me to model the 2027 cohort impact?"),
        ('"Which crews are behind on Jima quality this week?"',
         "Crew C: 61% A-Grade vs. the 80% target across 847 cuts yesterday. Crew Leader: Rodrigo Vázquez. Flag to his supervisor?"),
        ('"Are we on track for the Q3 distillery commitment?"',
         "Committed: 1,840 t. Current forecast: 1,710 t — 7% short. Window closes in 9 weeks. Draft a revised schedule?"),
    ]
    top = Inches(2.65)
    for q, a in chats:
        # Question
        qb = sl.shapes.add_shape(1, Inches(0.7), top, Inches(5.8), Inches(0.45))
        qb.fill.solid(); qb.fill.fore_color.rgb = SURFACE2; qb.line.color.rgb = BORDER
        qt = qb.text_frame; qt.word_wrap = True
        qr = qt.paragraphs[0].add_run(); qr.text = q
        qr.font.size = Pt(10); qr.font.italic = True; qr.font.color.rgb = TEXT2; qr.font.name = "Calibri"
        top += Inches(0.5)
        # Answer
        ab = sl.shapes.add_shape(1, Inches(1.1), top, Inches(5.4), Inches(0.5))
        ab.fill.solid(); ab.fill.fore_color.rgb = SURFACE; ab.line.color.rgb = GREEN
        at = ab.text_frame; at.word_wrap = True
        ar = at.paragraphs[0].add_run(); ar.text = a
        ar.font.size = Pt(10); ar.font.color.rgb = TEXT; ar.font.name = "Calibri"
        top += Inches(0.58)


def content_addon_modules(sl):
    modules = [
        ("⚙  Kaizen Engine", "Autonomous Continual Improvement — hunts inefficiency in your execution data and surfaces ranked OPEX reduction opportunities with dollar impact estimates."),
        ("🧠  The Strategist", "Board-level scenario modeling: harvest timing trade-offs, supply contract risk, multi-cohort capital planning, COGS benchmarking."),
        ("🚜  Fleet Command", "AI-optimized machine routing, predictive maintenance, real-time telematics for every tractor, drone, and transport vehicle."),
        ("🌾  Harvest Command", "AI orchestration for harvest: crew sequencing, transport coordination, real-time yield auditing."),
        ("🔔  Alert Builder", "Define monitoring rules in plain English. Context-aware alerts with AI-generated explanations, routed to the right person."),
    ]
    top = Inches(1.85)
    for name, desc in modules:
        txb(sl, Inches(0.7), top, Inches(3.2), Inches(0.4),
            name, size=Pt(12), bold=True, color=TEXT)
        txb(sl, Inches(3.9), top, Inches(9.1), Inches(0.55),
            desc, size=Pt(11), color=TEXT2)
        top += Inches(0.72)
        line = sl.shapes.add_shape(1, Inches(0.7), top, Inches(12.2), Pt(1))
        line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
        top += Inches(0.06)


def content_compliance(sl):
    txb(sl, Inches(0.7), Inches(1.85), Inches(6.0), Inches(0.3),
        "Chemical Compliance", size=Pt(12), bold=True, color=TEXT)
    bullet_list(sl, Inches(0.7), Inches(2.2), Inches(6.0), Inches(1.5),
        ["Real-time depot inventory — lot numbers, expiry dates, authorized applications per crop.",
         "Pre-Harvest Interval enforcement: residue-leaving sprays blocked at prescription, not at export inspection.",
         "Automatic application records. GlobalG.A.P. / SQF audit packages on demand."])
    txb(sl, Inches(0.7), Inches(3.85), Inches(6.0), Inches(0.3),
        "Worker PPE Compliance", size=Pt(12), bold=True, color=TEXT)
    bullet_list(sl, Inches(0.7), Inches(4.2), Inches(6.0), Inches(1.3),
        ["Computer vision verifies correct PPE class before any high-hazard task unlocks in Field Agent.",
         "AI verifies the right PPE for the specific chemical — not just presence detection.",
         "Full photo-verified record for every worker, every task. Complete evidentiary defense."])
    roi1 = add_roi_box(sl, Inches(0.7), Inches(5.65), Inches(6.0), Inches(0.5), "", GOLD)
    r = roi1.paragraphs[0].add_run()
    r.text = "Failed MRL export test:  $80K–$250K+ per rejected container"
    r.font.size = Pt(11); r.font.color.rgb = GOLD; r.font.name = "Calibri"


def slide_pricing(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, VOID)
    txb(sl, Inches(0.7), Inches(0.35), Inches(12), Inches(0.7),
        "Three Tiers. One Platform.", size=Pt(28), bold=True, color=TEXT)
    add_rule(sl, Inches(0.7), Inches(1.1))
    tiers = [
        ("Starter", "FarmAgent Starter", SURFACE, TEXT,
         "Technology Core\nCentral Brain\nAgave Crop Module\nField Agent App",
         "Producers digitalizing for the first time"),
        ("Most Popular ⭐", "FarmAgent Pro", BLUE, WHITE,
         "Everything in Starter\nFarm Brain™\nPersonal Agent\nAlert Builder",
         "Established operations ready to optimize"),
        ("Enterprise", "FarmAgent Enterprise", SURFACE, TEXT,
         "Everything in Pro\nKaizen Engine\nThe Strategist\nFleet Command\nCompliance Suite",
         "Vertically integrated agribusinesses"),
    ]
    for i, (tier_label, name, bg, tc, features, target) in enumerate(tiers):
        x = Inches(0.7) + Inches(4.2) * i
        y_offset = Inches(0.15) if i == 1 else 0
        card = sl.shapes.add_shape(1, x, Inches(1.4) - y_offset, Inches(4.0), Inches(5.6))
        card.fill.solid(); card.fill.fore_color.rgb = bg
        card.line.color.rgb = BLUE if i == 1 else BORDER
        txb(sl, x + Inches(0.2), Inches(1.6) - y_offset, Inches(3.6), Inches(0.35),
            tier_label.upper(), size=Pt(9), color=GREY_L if i != 1 else RGBColor(0xFF,0xFF,0xFF))
        txb(sl, x + Inches(0.2), Inches(2.0) - y_offset, Inches(3.6), Inches(0.45),
            name, size=Pt(16), bold=True, color=tc)
        txb(sl, x + Inches(0.2), Inches(2.6) - y_offset, Inches(3.6), Inches(2.0),
            features, size=Pt(11), color=TEXT2 if i != 1 else RGBColor(0xFF,0xFF,0xFF))
        txb(sl, x + Inches(0.2), Inches(6.0) - y_offset, Inches(3.6), Inches(0.5),
            target, size=Pt(9.5), color=GREY_M if i != 1 else RGBColor(0xFF,0xFF,0xFF))
    txb(sl, Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
        "Per-hectare annual subscription · Pay for the hectares you enroll and the modules you activate",
        size=Pt(10), color=GREY_M, align=PP_ALIGN.CENTER)
    footer_bar(sl)


def slide_roi_numbers(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, VOID)
    txb(sl, Inches(0.7), Inches(0.35), Inches(12), Inches(0.7),
        "The Agave ROI Case — By the Numbers", size=Pt(28), bold=True, color=TEXT)
    add_rule(sl, Inches(0.7), Inches(1.1))
    rows = [
        ("Unrecovered revenue: 10% undetected mortality, 100 ha over 5 years", "$148,500", "Advanced Inventory", GOLD),
        ("Jima quality uplift: 55% → 82% A-Grade on 100 ha", "$68K–$113K / season", "Photographic QA", GOLD),
        ("Scouting labor reduction via targeted GPS routing", "40–60%", "Targeted Scouting", GREEN),
        ("Measurement speed: smartphone vs. physical calipers", "5× faster", "Visual Calipers", GREEN),
        ("Time from problem observed to crew dispatched", "Hours (vs. days)", "Visual Localizer", BLUE_B),
        ("Max liability: failed MRL export test per container", "$80K–$250K+", "Chemical Compliance", GOLD),
    ]
    top = Inches(1.5)
    for desc, val, source, col in rows:
        txb(sl, Inches(0.7), top, Inches(7.8), Inches(0.42), desc, size=Pt(12), color=TEXT2)
        txb(sl, Inches(8.6), top, Inches(2.5), Inches(0.42), val, size=Pt(14), bold=True, color=col)
        txb(sl, Inches(11.2), top, Inches(2.0), Inches(0.42), source, size=Pt(10), color=GREY_M)
        top += Inches(0.44)
        line = sl.shapes.add_shape(1, Inches(0.7), top, Inches(12.2), Pt(1))
        line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
        top += Inches(0.06)
    footer_bar(sl)


def slide_cta(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(sl, img("deck_aerial_bright.png"), brightness=0.42)
    # Headline
    box = sl.shapes.add_textbox(Inches(2.0), Inches(1.3), Inches(9.5), Inches(3.5))
    tf = box.text_frame; tf.word_wrap = True
    lines = [("The farms", TEXT), ("that can ", TEXT), ("see", GREEN), (" will outperform", TEXT), ("the ones that can't.", TEXT)]
    # Simpler single-paragraph approach
    first = True
    for text, col in [("The farms\nthat can ", TEXT), ("see", GREEN), ("\nwill outperform\nthe ones that can't.", TEXT)]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False if first else first
        r = p.add_run()
        r.text = text
        r.font.size = Pt(46); r.font.color.rgb = col; r.font.name = "Calibri"
    txb(sl, Inches(2.5), Inches(4.9), Inches(8.5), Inches(0.9),
        "FarmAgent is available now. Start with aerial inventory on your existing blocks — "
        "see your mortality rate, replanting window, and monthly production trajectory in the first 90 days.",
        size=Pt(13), color=TEXT2, align=PP_ALIGN.CENTER)
    # CTA Button
    btn = sl.shapes.add_shape(1, Inches(5.5), Inches(5.95), Inches(2.6), Inches(0.55))
    btn.fill.solid(); btn.fill.fore_color.rgb = BLUE; btn.line.fill.background()
    bt = btn.text_frame.paragraphs[0]
    bt.alignment = PP_ALIGN.CENTER
    br = bt.add_run(); br.text = "REQUEST A DEMO"
    br.font.size = Pt(11); br.font.bold = True; br.font.color.rgb = WHITE; br.font.name = "Calibri"
    txb(sl, Inches(2.0), Inches(6.7), Inches(9.5), Inches(0.4),
        "sales@fruitscout.ai  ·  fruitscout.ai/farmagent",
        size=Pt(11), color=GREY_M, align=PP_ALIGN.CENTER)
    footer_bar(sl)


# ── Main ─────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # 1 — Cover
    slide_cover(prs)

    # 2 — Section: The Problem
    slide_divider(prs, "01", "01 — The Problem",
                  [("Agave farming", TEXT), (" is a 7-year", TEXT), (" blind bet.", GREEN)],
                  bg_img="deck_aerial_dark.png", brightness=0.28)

    # 3 — Problem detail
    slide_two_col(prs, "The Problem",
                  "The Most Expensive Crop in the World — Still Managed on Paper",
                  content_problem, "deck_jima_cut.png")

    # 4 — Section: The Solution
    slide_divider(prs, "02", "02 — The Solution",
                  [("Now you can", TEXT), (" see everything.", GREEN)],
                  bg_img="deck_aerial_bright.png", brightness=0.45)

    # 5 — Overview
    slide_two_col(prs, "FarmAgent Platform",
                  "Sight Across the Entire Operation",
                  content_overview, "deck_aerial_bright.png", rule_color=GREEN)

    # 6 — Section: The Platform
    slide_divider(prs, "03", "03 — The Platform",
                  [("Eyes in the sky.", TEXT), (" Eyes in every hand.", GREEN)],
                  bg_img="deck_aerial_dark.png", brightness=0.28)

    # 7 — Aerial Intelligence
    slide_two_col(prs, "Technology Core · Aerial Intelligence",
                  "Know Exactly What You Have — From 120 Metres Up",
                  content_aerial, "deck_drone_survey.png")

    # 8 — Field Agent
    slide_two_col(prs, "Technology Core · Field Agent",
                  "The Guide That Walks Every Worker Through Their Day — and Sees What They See",
                  content_field_agent, "deck_muddy_road.png", rule_color=GREEN)

    # 9 — Precision Measurement
    slide_two_col(prs, "Technology Core · Visual Sensor Suite",
                  "Measurements, Not Estimates — From Every Smartphone on the Estate",
                  content_measurement, "deck_jima_cut.png")

    # 10 — Central Brain
    slide_full_content(prs, "Technology Core · Operating System",
                       "Total Operational Sight — Four Pillars, One System",
                       content_central_brain)

    # 11 — 7-Year Plan
    slide_full_content(prs, "Technology Core · Plan Management",
                       "The Entire Lifecycle — One Orchestrated Plan",
                       content_seven_year)

    # 12 — Jima Quality
    slide_two_col(prs, "Add-On Module · Photographic QA",
                  "Grade Every Cut. At Scale. In Real Time.",
                  content_jima, "deck_jima_cut.png", addon=True, rule_color=GREEN)

    # 13 — Farm Brain
    slide_two_col(prs, "Technology Core · Learning Engine",
                  "Farm Brain™ — The Operation Gets Smarter Every Season",
                  content_farm_brain, "deck_aerial_bright.png")

    # 14 — Section: Intelligence Modules
    slide_divider(prs, "04", "04 — Intelligence Modules",
                  [("Extend the sight to", TEXT), (" every corner", GREEN), (" of the operation.", TEXT)],
                  bg_img="deck_aerial_dark.png", brightness=0.28)

    # 15 — Personal Agent
    slide_two_col(prs, "Add-On Module · Personal Agent",
                  "You Talk to It Like You Talk to Any Other Teammate",
                  content_personal_agent, "deck_personal_agent.png", addon=True, rule_color=GREEN)

    # 16 — Add-On Modules
    slide_full_content(prs, "Add-On Modules",
                       "Specialist AI Agents — Activated When You're Ready",
                       content_addon_modules, addon=True)

    # 17 — Compliance
    slide_two_col(prs, "Add-On Module · Security & Compliance",
                  "One Failed Shipment Costs More Than the Entire Platform",
                  content_compliance, "deck_ppe_worker.png", addon=True)

    # 18 — Section: Pricing
    slide_divider(prs, "05", "05 — Packaging",
                  [("Pay for what", TEXT), (" you protect.", GREEN)],
                  bg_img="deck_jima_cut.png", brightness=0.30)

    # 19 — Pricing Tiers
    slide_pricing(prs)

    # 20 — ROI Numbers
    slide_roi_numbers(prs)

    # 21 — CTA
    slide_cta(prs)

    out = os.path.join(os.path.dirname(__file__), "..", "FarmAgent-Agave.pptx")
    prs.save(out)
    print(f"Saved: {os.path.abspath(out)}")


if __name__ == "__main__":
    build()
